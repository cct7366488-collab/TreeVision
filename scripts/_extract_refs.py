"""Extract text content from the 3 reference files into raw text for synthesis.

This is a one-shot scaffolding helper, not a permanent script.
Writes UTF-8 output files directly via Python, bypassing shell encoding.
"""

import os
import sys

BASE = r"H:\我的雲端硬碟\cct\林業技師事務所\標案\012林業保育署台中分署\114-115年林業及自然保育署臺中分署轄區林業永續多元輔導計畫"
OUT_DIR = r"C:\Users\cct\projects\TreeVision\.claude\worktrees\sad-mendel-c5b43d\.scratch"
os.makedirs(OUT_DIR, exist_ok=True)

DOCX = os.path.join(BASE, "土肉桂施肥與修枝矮化試驗樣區設計報告書(1).docx")
XLSX = os.path.join(BASE, "土肉桂生長調查記錄表.xlsx")
PPTX = os.path.join(BASE, "土肉桂矮化修剪技術.pptx")


def extract_docx(path: str, out_path: str) -> None:
    import docx
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    doc = docx.Document(path)

    def iter_block_items(parent):
        if isinstance(parent, _Doc):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"# DOCX: {os.path.basename(path)}\n\n")
        block_idx = 0
        for block in iter_block_items(doc):
            block_idx += 1
            if isinstance(block, Paragraph):
                txt = block.text.strip()
                if txt:
                    style = block.style.name if block.style else ""
                    f.write(f"[P #{block_idx}][{style}] {txt}\n")
            elif isinstance(block, Table):
                f.write(f"\n[TABLE #{block_idx}] rows={len(block.rows)}, cols={len(block.columns)}\n")
                for r_i, row in enumerate(block.rows):
                    cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
                    f.write(f"  R{r_i}: " + " ┃ ".join(cells) + "\n")
                f.write("\n")


def extract_xlsx(path: str, out_path: str) -> None:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"# XLSX: {os.path.basename(path)}\n\n")
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            f.write(f"\n## Sheet: {sheet_name}  (dim={ws.dimensions}, max_row={ws.max_row}, max_col={ws.max_column})\n\n")
            max_r = min(ws.max_row, 200)
            max_c = min(ws.max_column, 30)
            for r in range(1, max_r + 1):
                row_vals = []
                for c in range(1, max_c + 1):
                    v = ws.cell(row=r, column=c).value
                    row_vals.append("" if v is None else str(v).strip())
                if any(row_vals):
                    f.write(f"R{r:>3}: " + " | ".join(row_vals) + "\n")
            if ws.max_row > max_r:
                f.write(f"... (truncated; total {ws.max_row} rows)\n")
        # Also dump merged cell info on first sheet
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            if ws.merged_cells.ranges:
                f.write(f"\n## Merged ranges in {sheet_name}\n")
                for mr in ws.merged_cells.ranges:
                    f.write(f"  {mr}\n")


def extract_pptx(path: str, out_path: str) -> None:
    from pptx import Presentation

    prs = Presentation(path)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"# PPTX: {os.path.basename(path)}\n\n")
        f.write(f"Total slides: {len(prs.slides)}\n\n")
        for i, slide in enumerate(prs.slides, start=1):
            f.write(f"\n## Slide {i}\n\n")
            title_shape = slide.shapes.title if slide.shapes.title else None
            if title_shape is not None:
                title_text = (title_shape.text or "").strip()
                if title_text:
                    f.write(f"**TITLE:** {title_text}\n\n")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == title_shape:
                    continue
                for para in shape.text_frame.paragraphs:
                    txt = (para.text or "").strip()
                    if txt:
                        f.write(f"- {txt}\n")
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    f.write(f"\n_[NOTES]_ {notes}\n")


if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else "all"
    if target in ("all", "docx"):
        extract_docx(DOCX, os.path.join(OUT_DIR, "ref_docx.txt"))
        print(f"DOCX extracted -> {os.path.join(OUT_DIR, 'ref_docx.txt')}")
    if target in ("all", "xlsx"):
        extract_xlsx(XLSX, os.path.join(OUT_DIR, "ref_xlsx.txt"))
        print(f"XLSX extracted -> {os.path.join(OUT_DIR, 'ref_xlsx.txt')}")
    if target in ("all", "pptx"):
        extract_pptx(PPTX, os.path.join(OUT_DIR, "ref_pptx.txt"))
        print(f"PPTX extracted -> {os.path.join(OUT_DIR, 'ref_pptx.txt')}")
